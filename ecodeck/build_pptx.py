#!/usr/bin/env python3
"""
build_pptx.py — Convert EcoDialysis decks (content_en / content_fr) to .pptx.

True 16:9 (13.333 in × 7.5 in), one slide per entry in SLIDES_EN / SLIDES_FR
(53 slides per language: 35 main + 18 appendix).

HTML in `body` is flattened to readable, structured text:
- <h2>/<h3> become bold lines
- <p>, <li>, <td>/<th> become paragraphs / bullets / cells
- <table> becomes a real PPTX table
- <a class="src">…</a> chips are inlined as [Label]
- everything else is stripped to text

Two output files (alongside the HTML decks):
  /home/victor/_aprojects/EcoDialysis/EcoDialysis_Final_Pitch_EN.pptx
  /home/victor/_aprojects/EcoDialysis/EcoDialysis_Final_Pitch_FR.pptx
"""

import os, re, sys, html
from html.parser import HTMLParser

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from content_en import SLIDES_EN
from content_fr import SLIDES_FR

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# -------- design tokens --------
NAVY   = RGBColor(0x0A, 0x2A, 0x55)
NAVY2  = RGBColor(0x0E, 0x3C, 0x7A)
BLUE   = RGBColor(0x2E, 0x7F, 0xE8)
INK    = RGBColor(0x0B, 0x1A, 0x2E)
GREY   = RGBColor(0x56, 0x69, 0x8A)
LINE   = RGBColor(0xD9, 0xE3, 0xEF)
SKY    = RGBColor(0xE6, 0xF0, 0xFB)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.5)

# -------- HTML → structured blocks --------
class Block:
    __slots__ = ("kind", "text", "level", "rows")
    def __init__(self, kind, text="", level=0, rows=None):
        self.kind = kind        # "h2","h3","p","li","table","spacer"
        self.text = text
        self.level = level
        self.rows = rows or []  # list[list[str]] for tables

class Flattener(HTMLParser):
    """Flatten body HTML to a list of Block objects."""
    def __init__(self):
        super().__init__(convert_charrefs=True)
        self.blocks = []
        self.cur = None        # currently-open text block
        self.stack = []
        # table state
        self.in_table = False
        self.tbl_rows = []
        self.cur_row = None
        self.cell_buf = None
        self.cur_cell_is_header = False
        # inline modifiers
        self.bold_depth = 0
        self.in_chip = False

    def _open(self, kind, level=0):
        self._close()
        self.cur = Block(kind, "", level)

    def _close(self):
        if self.cur is not None:
            t = re.sub(r"\s+", " ", self.cur.text).strip()
            if t:
                self.cur.text = t
                self.blocks.append(self.cur)
            self.cur = None

    def _append(self, s):
        if self.cell_buf is not None:
            self.cell_buf += s
        elif self.cur is not None:
            self.cur.text += s
        # else dropped (text outside any block)

    def handle_starttag(self, tag, attrs):
        a = dict(attrs)
        if tag in ("h1","h2"):
            self._open("h2")
        elif tag == "h3":
            self._open("h3")
        elif tag == "p":
            cls = a.get("class","")
            self._open("p")
            if "lead" in cls:
                self.cur.level = 1  # lead style
        elif tag == "li":
            depth = sum(1 for s in self.stack if s in ("ul","ol"))
            self._open("li", level=max(0, depth-1))
        elif tag in ("ul","ol"):
            self.stack.append(tag)
        elif tag == "br":
            self._append("\n")
        elif tag in ("b","strong"):
            self.bold_depth += 1
        elif tag == "table":
            self._close()
            self.in_table = True
            self.tbl_rows = []
            self.stack.append("table")
        elif tag == "tr" and self.in_table:
            self.cur_row = []
        elif tag in ("td","th") and self.in_table:
            self.cell_buf = ""
            self.cur_cell_is_header = (tag == "th")
            if tag == "th":
                self.bold_depth += 1
        elif tag == "a":
            cls = a.get("class","")
            if "src" in cls:
                self.in_chip = True
                self._append("[")
        elif tag == "div":
            cls = a.get("class","")
            if "card" in cls or "col" in cls:
                # treat card boundary as a soft separator
                self._close()
        elif tag == "span":
            pass
        # everything else: ignore tag

    def handle_endtag(self, tag):
        if tag in ("h1","h2","h3","p","li"):
            self._close()
        elif tag in ("ul","ol"):
            if self.stack and self.stack[-1] == tag:
                self.stack.pop()
        elif tag in ("b","strong"):
            self.bold_depth = max(0, self.bold_depth - 1)
        elif tag == "table" and self.in_table:
            self.in_table = False
            if self.stack and self.stack[-1] == "table":
                self.stack.pop()
            self.blocks.append(Block("table", rows=self.tbl_rows))
            self.tbl_rows = []
        elif tag == "tr" and self.cur_row is not None:
            if self.cur_row:
                self.tbl_rows.append(self.cur_row)
            self.cur_row = None
        elif tag in ("td","th") and self.cell_buf is not None:
            txt = re.sub(r"\s+", " ", self.cell_buf).strip()
            self.cur_row.append(txt)
            self.cell_buf = None
            if tag == "th":
                self.bold_depth = max(0, self.bold_depth - 1)
        elif tag == "a" and self.in_chip:
            self.in_chip = False
            self._append("]")

    def handle_data(self, data):
        self._append(data)

    def close(self):
        self._close()
        return self.blocks

def flatten_html(html_str: str):
    f = Flattener()
    f.feed(html_str)
    return f.close()

# -------- PPTX rendering --------
def add_textbox(slide, left, top, width, height, blocks, body_font=14):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    first = True
    for b in blocks:
        if b.kind == "table":
            continue  # handled separately
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if b.kind == "h2":
            p.text = b.text
            for r in p.runs:
                r.font.bold = True; r.font.size = Pt(18); r.font.color.rgb = NAVY2
            p.space_after = Pt(4)
        elif b.kind == "h3":
            p.text = b.text
            for r in p.runs:
                r.font.bold = True; r.font.size = Pt(14); r.font.color.rgb = NAVY2
            p.space_after = Pt(2)
        elif b.kind == "li":
            p.text = "• " + b.text
            p.level = b.level
            for r in p.runs:
                r.font.size = Pt(body_font); r.font.color.rgb = INK
            p.space_after = Pt(2)
        else:  # p, spacer, fallback
            p.text = b.text
            size = body_font + (2 if b.level == 1 else 0)
            for r in p.runs:
                r.font.size = Pt(size); r.font.color.rgb = INK
                if b.level == 1:
                    r.font.color.rgb = NAVY
            p.space_after = Pt(4)
    return tb

def add_table(slide, left, top, width, height, rows):
    if not rows:
        return None
    ncols = max(len(r) for r in rows)
    rows = [r + [""]*(ncols-len(r)) for r in rows]
    nrows = len(rows)
    # cap dimensions
    tbl_shape = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    tbl = tbl_shape.table
    for ci in range(ncols):
        tbl.columns[ci].width = width // ncols
    for ri, row in enumerate(rows):
        is_header = (ri == 0)
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Emu(36000); tf.margin_right = Emu(36000)
            tf.margin_top = Emu(18000); tf.margin_bottom = Emu(18000)
            p = tf.paragraphs[0]
            p.text = cell_text
            for r in p.runs:
                r.font.size = Pt(10)
                r.font.color.rgb = WHITE if is_header else INK
                r.font.bold = is_header
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if is_header else (SKY if ri % 2 == 0 else WHITE)
    return tbl_shape

def add_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg

def add_accent_bar(slide, color=BLUE, y=Inches(0.85)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, y, Inches(1.2), Emu(45000))
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()

def add_title(slide, text, color=NAVY, size=30):
    tb = slide.shapes.add_textbox(MARGIN, Inches(0.3), SLIDE_W - 2*MARGIN, Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_top = Emu(0)
    p = tf.paragraphs[0]; p.text = text
    for r in p.runs:
        r.font.bold = True; r.font.size = Pt(size); r.font.color.rgb = color
    return tb

def add_footer(slide, brand, part, counter):
    y = SLIDE_H - Inches(0.4)
    for x, w, txt, align, col in [
        (MARGIN, Inches(5.5), brand, PP_ALIGN.LEFT, GREY),
        (Inches(5.5), Inches(2.3), part, PP_ALIGN.CENTER, GREY),
        (SLIDE_W - MARGIN - Inches(2.0), Inches(2.0), counter, PP_ALIGN.RIGHT, GREY),
    ]:
        tb = slide.shapes.add_textbox(x, y, w, Inches(0.3))
        tf = tb.text_frame
        tf.margin_left = Emu(0); tf.margin_top = Emu(0)
        p = tf.paragraphs[0]; p.text = txt; p.alignment = align
        for r in p.runs:
            r.font.size = Pt(10); r.font.color.rgb = col

def render_slide(prs, idx, total, sid, kind, title, body_html, notes, brand, part_map_get):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_bg(s, WHITE)
    part = part_map_get(sid, kind)

    if kind == "cover":
        # full navy cover
        add_bg(s, NAVY)
        tb = s.shapes.add_textbox(MARGIN, Inches(2.3), SLIDE_W - 2*MARGIN, Inches(2.0))
        p = tb.text_frame.paragraphs[0]; p.text = title; p.alignment = PP_ALIGN.LEFT
        for r in p.runs:
            r.font.bold = True; r.font.size = Pt(54); r.font.color.rgb = WHITE
        # subtitle from body (strip first paragraph if any)
        blocks = flatten_html(body_html)
        sub = next((b.text for b in blocks if b.kind in ("p","h2","h3")), "")
        if sub:
            tb2 = s.shapes.add_textbox(MARGIN, Inches(4.1), SLIDE_W - 2*MARGIN, Inches(1.5))
            tf = tb2.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = sub
            for r in p.runs:
                r.font.size = Pt(20); r.font.color.rgb = SKY
        # brand stripe
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(2.0), Inches(1.2), Emu(45000))
        bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
        # footer in white
        ftb = s.shapes.add_textbox(MARGIN, SLIDE_H - Inches(0.5), SLIDE_W - 2*MARGIN, Inches(0.3))
        p = ftb.text_frame.paragraphs[0]; p.text = brand
        for r in p.runs:
            r.font.size = Pt(11); r.font.color.rgb = SKY
        _set_notes(s, notes)
        return

    if kind == "divider":
        add_bg(s, NAVY2)
        tb = s.shapes.add_textbox(MARGIN, Inches(2.7), SLIDE_W - 2*MARGIN, Inches(2.0))
        p = tb.text_frame.paragraphs[0]; p.text = title; p.alignment = PP_ALIGN.LEFT
        for r in p.runs:
            r.font.bold = True; r.font.size = Pt(40); r.font.color.rgb = WHITE
        blocks = flatten_html(body_html)
        sub = next((b.text for b in blocks if b.kind in ("p","h2","h3")), "")
        if sub:
            tb2 = s.shapes.add_textbox(MARGIN, Inches(4.4), SLIDE_W - 2*MARGIN, Inches(1.5))
            p = tb2.text_frame.paragraphs[0]; p.text = sub
            for r in p.runs:
                r.font.size = Pt(18); r.font.color.rgb = SKY
        ftb = s.shapes.add_textbox(MARGIN, SLIDE_H - Inches(0.5), SLIDE_W - 2*MARGIN, Inches(0.3))
        p = ftb.text_frame.paragraphs[0]; p.text = f"{brand}    ·    {part}    ·    {idx} / {total}"
        for r in p.runs:
            r.font.size = Pt(10); r.font.color.rgb = SKY
        _set_notes(s, notes)
        return

    # content / appendix
    add_title(s, title, color=NAVY, size=24 if len(title) < 80 else 20)
    add_accent_bar(s, BLUE)
    blocks = flatten_html(body_html)
    text_blocks = [b for b in blocks if b.kind != "table"]
    table_blocks = [b for b in blocks if b.kind == "table"]

    top = Inches(1.05)
    avail_w = SLIDE_W - 2*MARGIN
    avail_h = SLIDE_H - top - Inches(0.6)

    if table_blocks and text_blocks:
        # text on top half, first table below
        add_textbox(s, MARGIN, top, avail_w, avail_h * 0.45, text_blocks, body_font=12)
        add_table(s, MARGIN, top + avail_h * 0.5, avail_w, avail_h * 0.5, table_blocks[0].rows)
    elif table_blocks:
        add_table(s, MARGIN, top, avail_w, avail_h, table_blocks[0].rows)
    else:
        # pick font size based on text load
        n_chars = sum(len(b.text) for b in text_blocks)
        fs = 14 if n_chars < 1200 else (12 if n_chars < 2200 else 10)
        add_textbox(s, MARGIN, top, avail_w, avail_h, text_blocks, body_font=fs)

    add_footer(s, brand, part, f"{idx} / {total}")
    _set_notes(s, notes)

def _set_notes(slide, notes):
    if not notes:
        return
    ns = slide.notes_slide
    tf = ns.notes_text_frame
    tf.text = notes

# -------- part labels (mirror build.py) --------
PART_MAP_EN = {
  "title":"Cover","agenda":"Agenda","why-now":"Part 1 · Hook","why-dialysis":"Part 1 · Hook","what-wasted":"Part 1 · Hook",
  "div-measure":"Part 2","measure":"Part 2 · Measurement","reference":"Part 2 · Reference site","reuse-map":"Part 2 · Reuse map",
  "legal-eich":"Part 2 · Regulation","dossier":"Part 2 · Dossier","tailored":"Part 2 · Tailored fit",
  "div-proof":"Part 3","intl-map":"Part 3 · Proof","uk-cluster":"Part 3 · UK","anz-cluster":"Part 3 · ANZ",
  "optim":"Part 3 · Optimisation","fr-proof":"Part 3 · France",
  "div-product":"Part 4","not-tank":"Part 4 · Product","product":"Part 4 · Product","ecosystem":"Part 4 · Equipment",
  "maturity":"Part 4 · Maturity","business":"Part 4 · Business","buyers":"Part 4 · Buyers",
  "div-ask":"Part 5","pilots":"Part 5 · Pilots","need":"Part 5 · Need","deliver":"Part 5 · Delivery",
  "refuse":"Part 5 · Discipline","kpis":"Part 5 · KPIs","risk":"Part 5 · Risk",
  "defensibility":"Part 5 · Moat","expansion":"Part 5 · Expansion","close":"Close",
}
PART_MAP_FR = {
  "title":"Couverture","agenda":"Sommaire","why-now":"Partie 1 · Accroche","why-dialysis":"Partie 1 · Accroche","what-wasted":"Partie 1 · Accroche",
  "div-measure":"Partie 2","measure":"Partie 2 · Mesure","reference":"Partie 2 · Site de référence","reuse-map":"Partie 2 · Carte réutilisation",
  "legal-eich":"Partie 2 · Réglementation","dossier":"Partie 2 · Dossier","tailored":"Partie 2 · Sur mesure",
  "div-proof":"Partie 3","intl-map":"Partie 3 · Preuves","uk-cluster":"Partie 3 · UK","anz-cluster":"Partie 3 · ANZ",
  "optim":"Partie 3 · Optimisation","fr-proof":"Partie 3 · France",
  "div-product":"Partie 4","not-tank":"Partie 4 · Produit","product":"Partie 4 · Produit","ecosystem":"Partie 4 · Équipement",
  "maturity":"Partie 4 · Maturité","business":"Partie 4 · Modèle","buyers":"Partie 4 · Acheteurs",
  "div-ask":"Partie 5","pilots":"Partie 5 · Pilotes","need":"Partie 5 · Besoin","deliver":"Partie 5 · Livraison",
  "refuse":"Partie 5 · Discipline","kpis":"Partie 5 · KPIs","risk":"Partie 5 · Risque",
  "defensibility":"Partie 5 · Défense","expansion":"Partie 5 · Expansion","close":"Clôture",
}

def build(slides, lang, brand, part_map, out_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    total = len(slides)
    def part_for(sid, kind):
        if kind == "appendix" or sid.startswith("apx-") or sid == "div-appendix":
            return "Appendix" if lang == "en" else "Annexe"
        return part_map.get(sid, "")
    for i, (sid, kind, title, body, notes) in enumerate(slides):
        render_slide(prs, i+1, total, sid, kind, title, body, notes, brand, part_for)
    prs.save(out_path)
    return out_path

if __name__ == "__main__":
    OUT = "/home/victor/_aprojects/EcoDialysis"
    brand = "EcoDialysis · ANDCS × MIT Hacking Medicine · Paris 2026"
    p1 = build(SLIDES_EN, "en", brand, PART_MAP_EN, f"{OUT}/EcoDialysis_Final_Pitch_EN.pptx")
    p2 = build(SLIDES_FR, "fr", brand, PART_MAP_FR, f"{OUT}/EcoDialysis_Final_Pitch_FR.pptx")
    print("EN:", p1, os.path.getsize(p1), "bytes")
    print("FR:", p2, os.path.getsize(p2), "bytes")
